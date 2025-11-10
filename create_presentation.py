"""
Create PowerPoint Presentation for CCE 313 - ICMP Pinger Lab
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors - Standard Professional Theme
    bg_color = RGBColor(255, 255, 255)  # White background
    accent_color = RGBColor(68, 114, 196)  # Professional blue
    text_color = RGBColor(0, 0, 0)  # Black text
    subtitle_color = RGBColor(89, 89, 89)  # Gray subtitle
    
    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = accent_color
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(20)
        subtitle_frame.paragraphs[0].font.color.rgb = subtitle_color
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, content_items):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = accent_color
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.7))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for item in content_items:
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = text_color
            p.space_before = Pt(8)
            p.level = 0 if not item.startswith('  ') else 1
        
        return slide
    
    # Slide 1: Title Slide
    add_title_slide(
        "CCE 313 - Computer Networks",
        "ICMP Pinger Lab: Network Diagnostics Tool\n\nProsenjit Mondol\nID: 2102049 | Reg: 10176\nB.Sc. CSE, 5th Semester\nNovember 10, 2025"
    )
    
    # Slide 2: Introduction & Problem Statement
    add_content_slide(
        "Introduction & Problem Statement",
        [
            "INTRODUCTION",
            "• Network performance monitoring is critical for reliable communication",
            "• ICMP (Internet Control Message Protocol) - fundamental for diagnostics",
            "",
            "PROBLEM STATEMENT",
            "Network administrators need tools to:",
            "  • Measure network latency and round-trip time (RTT)",
            "  • Detect packet loss and network congestion",
            "  • Monitor network stability and jitter",
            "  • Visualize real-time network performance",
            "",
            "IMPORTANCE IN NETWORKING",
            "• Quality of Service (QoS) monitoring for SLA compliance",
            "• Network troubleshooting and performance optimization",
            "• Real-time diagnostics for instant network health feedback"
        ]
    )
    
    # Slide 3: ICMP Protocol Architecture
    add_content_slide(
        "ICMP Protocol Architecture (RFC 792)",
        [
            "ICMP OVERVIEW - Network Layer (Layer 3)",
            "• Error Reporting: Unreachable destinations, TTL exceeded",
            "• Diagnostic Functions: Echo Request/Reply (Ping)",
            "• Network Control Messages: Redirect, Source Quench",
            "",
            "ICMP PACKET STRUCTURE (Echo Request/Reply)",
            "┌─────────────────────────────────┐",
            "│ IP Header (20 bytes)            │",
            "├─────────────────────────────────┤",
            "│ Type (8 bits) - Request:8/Reply:0│",
            "├─────────────────────────────────┤",
            "│ Code (8 bits) - 0               │",
            "├─────────────────────────────────┤",
            "│ Checksum (16 bits)              │",
            "├─────────────────────────────────┤",
            "│ Identifier (16 bits)            │",
            "├─────────────────────────────────┤",
            "│ Sequence Number (16 bits)       │",
            "├─────────────────────────────────┤",
            "│ Data (Variable)                 │",
            "└─────────────────────────────────┘"
        ]
    )
    
    # Slide 4: Implementation Methodology
    add_content_slide(
        "Implementation Methodology",
        [
            "THREE-TIER IMPLEMENTATION APPROACH",
            "",
            "1. RAW SOCKET IMPLEMENTATION (Scapy)",
            "  • Full control over ICMP packet structure",
            "  • Requires administrator privileges",
            "  • packet = IP(dst=host) / ICMP(type=8, code=0)",
            "",
            "2. WINDOWS ICMP API (No Admin Required)",
            "  • Windows IP Helper API (iphlpapi.dll)",
            "  • IcmpSendEcho function for echo requests",
            "  • Works without elevated privileges",
            "",
            "3. SUBPROCESS FALLBACK (Cross-Platform)",
            "  • Platform-native ping commands",
            "  • Maximum compatibility across OS",
            "  • Regex parsing for RTT extraction",
            "",
            "ARCHITECTURE: User Input → ICMPPinger → Implementation",
            "Selection → Results Processing → GUI Update"
        ]
    )
    
    # Slide 5: Advanced Features
    add_content_slide(
        "Advanced Features & Components",
        [
            "REAL-TIME PERFORMANCE MONITORING",
            "• Live RTT Graphing (Matplotlib Integration)",
            "  - Dynamic plotting of latency trends",
            "  - Color-coded success/failure indicators",
            "",
            "COMPREHENSIVE STATISTICS DASHBOARD",
            "• Minimum/Maximum/Average RTT",
            "• Jitter (network stability metric)",
            "• Packet Loss % and Success Rate %",
            "",
            "MULTI-TAB GUI INTERFACE (Tkinter)",
            "• Ping Client Tab: Configuration and execution",
            "• Echo Server Tab: Simulated server responses",
            "• Statistics Tab: Detailed analytics and CSV export",
            "",
            "KEY FEATURES",
            "✓ No Admin Required (Windows API fallback)",
            "✓ Cross-Platform Support (Windows/Linux/Mac)",
            "✓ Thread-Safe Operation | ✓ Export Capabilities"
        ]
    )
    
    # Slide 6: Analysis & Results
    add_content_slide(
        "Analysis & Results",
        [
            "TEST SCENARIO",
            "Target: 8.8.8.8 (Google DNS) | Count: 20 | Interval: 1s | Timeout: 10s",
            "",
            "SAMPLE RESULTS",
            "─────────────────────────────",
            "Minimum RTT:      12.34 ms",
            "Maximum RTT:      45.67 ms",
            "Average RTT:      23.45 ms",
            "Jitter:            8.92 ms",
            "Packet Loss:       0.0%",
            "Success Rate:    100.0%",
            "",
            "ANALYSIS INSIGHTS",
            "• Low Packet Loss (0%) → Stable network connection",
            "• Moderate Jitter (8.92ms) → Acceptable variation",
            "• Average RTT (23.45ms) → Good connection quality",
            "",
            "REAL-WORLD IMPLICATIONS",
            "• Network Health Assessment for SLA compliance",
            "• Troubleshooting congestion | Route optimization",
            "• Application Performance (VoIP/Gaming require <50ms)"
        ]
    )
    
    # Slide 7: Code Implementation
    add_content_slide(
        "Practical Implementation",
        [
            "1. SCAPY IMPLEMENTATION",
            "packet = IP(dst=host) / ICMP(seq=seq_num)",
            "response = sr1(packet, timeout=10.0)",
            "if response.getlayer(ICMP).type == 0:",
            "    rtt_ms = (end_time - start_time) * 1000",
            "",
            "2. WINDOWS ICMP API",
            "handle = iphlpapi.IcmpCreateFile()",
            "iphlpapi.IcmpSendEcho(handle, addr, data,",
            "    size, None, reply_buf, buf_size, timeout)",
            "reply = ICMP_ECHO_REPLY.from_buffer(reply_buf)",
            "rtt = reply.RoundTripTime",
            "",
            "3. SUBPROCESS FALLBACK",
            "cmd = ['ping', '-n', '1', '-w', '10000', host]",
            "output = subprocess.run(cmd, capture_output=True)",
            "match = re.search(r'time[=<]\\s*([\\d\\.]+)\\s*ms', output)",
            "rtt = float(match.group(1))"
        ]
    )
    
    # Slide 8: Performance Metrics
    add_content_slide(
        "Network Performance Metrics",
        [
            "STATISTICAL CALCULATIONS",
            "",
            "1. Round-Trip Time (RTT) Analysis",
            "   • Min, Max, Average RTT calculations",
            "   • RTT_avg = (1/n) × Σ RTT_i",
            "",
            "2. Jitter Calculation (Network Stability)",
            "   • Jitter = (1/n-1) × Σ |RTT_(i+1) - RTT_i|",
            "",
            "3. Packet Loss Percentage",
            "   • Packet Loss = (Failed / Total) × 100%",
            "",
            "PERFORMANCE INTERPRETATION",
            "RTT:    Excellent <20ms | Good 20-50ms | Fair 50-100ms",
            "Jitter: Excellent <5ms  | Good 5-20ms  | Fair 20-50ms",
            "Loss:   Excellent 0%    | Good <1%     | Fair 1-5%",
            "",
            "EXPORT FORMAT: CSV with metrics, values, and units"
        ]
    )
    
    # Slide 9: Challenges & Solutions
    add_content_slide(
        "Challenges & Solutions",
        [
            "TECHNICAL CHALLENGES ENCOUNTERED",
            "",
            "1. Permission Issues (Raw Sockets)",
            "   Challenge: Raw socket creation requires admin privileges",
            "   Solution: Implemented Windows ICMP API fallback",
            "   Result: Works without administrator rights on Windows",
            "",
            "2. Cross-Platform Compatibility",
            "   Challenge: Different OS-specific implementations",
            "   Solution: Multi-tier fallback (Scapy → API → Subprocess)",
            "   Result: Seamless operation across Windows/Linux/macOS",
            "",
            "3. Thread Safety in GUI Updates",
            "   Challenge: Matplotlib updates causing conflicts",
            "   Solution: Used canvas.draw_idle() and Tkinter.after()",
            "   Result: Smooth real-time graph updates",
            "",
            "BEST PRACTICES",
            "✓ Graceful degradation | ✓ Comprehensive error handling",
            "✓ Resource management  | ✓ Real-time user feedback"
        ]
    )
    
    # Slide 10: Conclusion
    add_content_slide(
        "Conclusion & Summary",
        [
            "KEY ACHIEVEMENTS",
            "• Implemented three approaches for ICMP echo requests",
            "• Developed GUI with live RTT graphing and statistics",
            "• Created cross-platform solution (Windows/Linux/macOS)",
            "• Built practical network diagnostic tool",
            "",
            "MAIN TAKEAWAYS",
            "• ICMP is essential for network layer diagnostics",
            "• Multi-tier implementation ensures reliability",
            "• Statistical analysis provides actionable insights",
            "• GUI-based tools make diagnostics accessible",
            "",
            "PROJECT OUTCOMES",
            "✓ Functional ICMP ping client (3 methods)",
            "✓ Real-time RTT visualization",
            "✓ Comprehensive statistics (min/max/avg/jitter/loss)",
            "✓ Export capabilities (CSV reports)",
            "",
            "FUTURE ENHANCEMENTS",
            "• Traceroute integration | • Wireshark integration",
            "• Multiple host monitoring | • Historical data trending"
        ]
    )
    
    # Slide 11: References
    add_content_slide(
        "References",
        [
            "TECHNICAL DOCUMENTATION",
            "• RFC 792 - Internet Control Message Protocol",
            "  https://tools.ietf.org/html/rfc792",
            "",
            "• Forouzan, B. A. (2012)",
            "  Data Communications and Networking, 5th Edition",
            "  McGraw-Hill Education",
            "",
            "• Kurose, J. F., & Ross, K. W. (2021)",
            "  Computer Networking: A Top-Down Approach, 8th Edition",
            "  Pearson Education",
            "",
            "IMPLEMENTATION LIBRARIES",
            "• Scapy Documentation: https://scapy.readthedocs.io/",
            "• Python ctypes Module",
            "• Windows IP Helper API (Microsoft Docs)",
            "",
            "PROJECT RESOURCES",
            "GitHub: https://github.com/pronad1/Socket_Programming",
            "Email: prosenjit1156@gmail.com"
        ]
    )
    
    # Thank You Slide
    add_title_slide(
        "Thank You! 🎓",
        "Questions & Discussion\n\nProsenjit Mondol\nID: 2102049 | Reg: 10176\n\nContact: prosenjit1156@gmail.com\nGitHub: https://github.com/pronad1"
    )
    
    # Save presentation
    output_file = "CCE313_ICMP_Pinger_Presentation.pptx"
    prs.save(output_file)
    print(f"✓ Presentation created successfully: {output_file}")
    print(f"✓ Total slides: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    create_presentation()
